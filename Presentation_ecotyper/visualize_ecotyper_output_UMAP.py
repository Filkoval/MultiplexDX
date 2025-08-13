import pandas as pd
import sys
import os
import matplotlib.pyplot as plt
from seaborn import color_palette
import subprocess
import shutil
import fitz 
from pptx import Presentation
from PIL import Image
from sklearn.cluster import DBSCAN
import hdbscan
import umap.plot
import scanpy as sc
import umap.umap_ as _um
import numpy as np
from sklearn.metrics import silhouette_score

UMAP = _um.UMAP

# ======= DATA LOADING ========
DIR = sys.argv[1]
ENTRIES = os.listdir(DIR)
ECOTYPER_DIRS = [d for d in ENTRIES if os.path.isdir(os.path.join(DIR,d)) and d.startswith("Sample-")]
SPACE_RANGER_DIRS = [d for d in ENTRIES if os.path.isdir(os.path.join(DIR,d)) and d.startswith("Space_ranger-")]
ANNOTATION_FILES = [f for f in ENTRIES if os.path.isfile(os.path.join(DIR,f)) and f.startswith("annot-")]
CLUSTERS = 4
ECOTYPES = ["CE"+str(i) for i in range(1,11)]
ATTEMPTS = 10
POINT_SIZES = [6,5,8,6]
PATH = os.getcwd()


# ======== CREATES OUTPUT DIRECTORY FOR PLOTS =========
def create_output_dir():
    path = os.path.join(DIR,"plots")
    if os.path.exists(path) and os.path.isdir(path):
        answear = input("Directory plots already exists, do you wish to override it? (y|n)")
        if answear == "y": 
            shutil.rmtree(os.path.join(DIR,"plots"))
            os.makedirs(path)
        elif answear == "n": 
            print("saving to /plots additionaly, warning! : may cause overlap of plots with the same name.")
            create_presentation()
            sys.exit()
        else: create_output_dir()
    else: os.makedirs(path)


# ======= CORRECTENSS OF THE INPUT =======
#  - input must contain at least as many ecotyper directories as cluster files
def input_ok():
    if len(SPACE_RANGER_DIRS) > len(ECOTYPER_DIRS): 
        missing_dirs = []
        for file in SPACE_RANGER_DIRS:
            filenames = file.split("Space_ranger-",1)
            if filenames[1] not in ECOTYPER_DIRS: missing_dirs.append(filenames[1]) 
        print("Incomplete input, ecotyper directories missing:")
        for filename in filenames:
            print(filename)
        return False
    if len(SPACE_RANGER_DIRS) < len(ECOTYPER_DIRS): 
        print("Missing space ranger directories for certain ecotyper directories, ommiting these samples")
    if len(ANNOTATION_FILES) > len(SPACE_RANGER_DIRS): 
        print("Missing space ranger directories for certain annotation images, ommiting these samples")
    if len(ANNOTATION_FILES) < len(SPACE_RANGER_DIRS): 
        print("Missing annotation images, please add images before proceeding")
        return False
    return True


# ======= EXTRACTS NAMES OF FILES ========
def extract(filename):
    return os.path.splitext(filename)[0]


# ======= DICTIONARY OF SAMPLE FILES ======= 
def create_dic():
    dic = {}
    filename = "gene_expression_kmeans_"+str(CLUSTERS)+"_clusters"
    for dir in ECOTYPER_DIRS:
        dirnames = dir.split("Sample-",1)
        dname = extract(dirnames[1])
        for space in SPACE_RANGER_DIRS:
            space_names = space.split("Space_ranger-",1)
            sname = extract(space_names[1])
            if dname == sname : 
                dic[dname] = [
                os.path.join(DIR,space,"analysis","clustering",filename,"clusters.csv"),
                os.path.join(os.path.join(DIR,dir),"ecotype_abundances.txt"),
                os.path.join(DIR,space)
                ]
    return dic


# ====== CLUSTERS BARCODES ======
def filter_barcode(cluster, df):
    barcodes = []
    for spot in df[df["Cluster"] == cluster].itertuples():
        barcodes.append(spot[1])
    return barcodes


# ====== PIE CHART ======
def draw_pie(values,name,cluster,path):
    palette = color_palette("pastel", len(values))
    filename = "piechart-"+name+"_"+str(cluster)+".png"
    pie_name = os.path.join(path,filename)
    plt.pie(
        values, 
        labels=ECOTYPES, 
        autopct='%1.1f%%', 
        startangle=90,
        colors=palette, 
        wedgeprops={'edgecolor': 'white', 'linewidth': 1.5},
        pctdistance=0.85
        )
    fig = plt.gcf()
    fig.savefig(pie_name, dpi=300, bbox_inches='tight')
    plt.close(fig)
    return pie_name


# ======= SPATIAL GRAPH ======== 
def draw_spatial(file_dic,name,cluster,path,point,cluster_csv):

    tissue_png = os.path.join(path,str(name+"_tissue_"+str(cluster)+".png"))
    tissue_pdf = os.path.join(path,str(name+"_tissue_"+str(cluster)+".pdf"))
    no_tissue_png = os.path.join(path,str(name+"_no_tissue_"+str(cluster)+".png"))
    no_tissue_pdf = os.path.join(path,str(name+"_no_tissue_"+str(cluster)+".pdf"))


    subprocess.run(
        [
            "Rscript",
            PATH+"\\spatial_plot.r",
            file_dic[sample][2],
            cluster_csv,
            #file_dic[sample][0],
            str(cluster),
            tissue_pdf,
            no_tissue_pdf,
            str(POINT_SIZES[point]),
            cluster_csv
        ],
        capture_output=True,    
        text=True,
        cwd=PATH, 
        ) 
    
    convert_to_png(tissue_pdf,tissue_png)
    convert_to_png(no_tissue_pdf,no_tissue_png,crop = True)


# ======= CONVERTS PDF TO PNG AND CROPS IT =========
def convert_to_png(path_in,path_out,crop = False):
    
    doc = fitz.open(path_in)
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=300)
        pix.save(path_out)
    doc.close()
    if crop: crop_px = 100
    else: crop_px = 30
    with Image.open(path_out) as img:
            width, height = img.size
            cropped = img.crop((
                crop_px,
                crop_px,
                width - crop_px,
                height - crop_px
            ))
            cropped.save(path_out)

# ======= CROP ANNOTATION ===========
def crop_annot(inimg,name):

    def extract_val(key):
        return int(inimg.split(f"{key}=",1)[1].split(",",1)[0])
    
    img = Image.open(os.path.join(DIR, inimg))
    csv_path = os.path.join(file_dic[name][2], "spatial", "tissue_positions.csv")
    csv = pd.read_csv(csv_path)

    width, height = img.size

    up = csv["array_row"].min()
    left = csv["array_col"].min()
    down = csv["array_row"].max()
    right = csv["array_col"].max()

    up_left = csv[(csv["array_row"] == up) & (csv["array_col"] == left)].iloc[0]
    down_right = csv[(csv["array_row"] == down) & (csv["array_col"] == right)].iloc[0]

    left_cut = extract_val("x")
    up_cut = extract_val("y")
    width_cut = extract_val("w")
    height_cut = int(inimg.split("h=",1)[1].split(")",1)[0])

    scale_x = width / width_cut
    scale_y = height / height_cut

    x_offset = left_cut * scale_x
    y_offset = up_cut * scale_y

    x1 = down_right["pxl_col_in_fullres"] * scale_x
    y1 = down_right["pxl_row_in_fullres"] * scale_y
    x2 = up_left["pxl_col_in_fullres"] * scale_x
    y2 = up_left["pxl_row_in_fullres"] * scale_y

    w, h = round(abs(x2 - x1)) + 1, round(abs(y2 - y1)) + 1
    ox = round(min(x1, x2) - x_offset)
    oy = round(min(y1, y2) - y_offset)

    left   = max(0, ox)
    top    = max(0, oy)
    right  = min(img.width, ox + w)
    bottom = min(img.height, oy + h)
    cropped = img.crop((left, top, right, bottom))

    out_path = os.path.join(DIR, "plots", name, "annotation.png")
    cropped.save(out_path)


# ======= CREATE SET OF GRAPHS ==========
def create_graphs(name, cluster,point,cluster_csv,values):
    dirname = "sample_" + name + "_cluster_" + str(cluster)
    path = os.path.join(DIR,"plots",sample,dirname)
    os.makedirs(path)
    pie = draw_pie(values,sample,i+1,path)
    annot_file = [f for f in ANNOTATION_FILES if f.startswith("annot-"+name)][0]
    if len(annot_file) == 0: return False
    draw_spatial(file_dic,name,cluster,path,point,cluster_csv)
    crop_annot(annot_file, name)
    return True


# ====== WRITE IN PRESENTATION =========
def write_text(ph_num,text,slide):
        title_placeholder = slide.placeholders[ph_num]
        title_placeholder.text = text

def add_img(ph_num,pic,slide):
    placeholder = slide.placeholders[ph_num]
    placeholder.insert_picture(pic)


# ======== CREATES POWERPOINT PRESENTATION ===========
def create_presentation():
    odd = 0
    prs = Presentation(os.path.join(DIR,"template.pptx"))
    samples = os.listdir(os.path.join(DIR,"plots"))
    placeholder_order = [18,17,19,14,13,15]
    sample_order = [16,12]

    for sample in samples:

        plots_root = os.path.join(DIR, "plots", sample)
        clusters = [
            name for name in os.listdir(plots_root)
            if os.path.isdir(os.path.join(plots_root, name))
        ]

        for cluster in clusters:
            imgs = [f for f in os.listdir(os.path.join(DIR,"plots",sample,cluster)) if not f.endswith("annot_1.png")and f.endswith(".png")]
            if odd == 0: 
                slide = prs.slides.add_slide(prs.slide_layouts[12])
                write_text(0,"Sample "+sample,slide)
                odd = 1
                for i in range(2):
                    add_img(sample_order[i],os.path.join(DIR,"plots",sample,"annotation.png"),slide)
            else: odd = 0
            for i, img in enumerate(imgs):
                add_img(placeholder_order[i+odd*3],os.path.join(DIR,"plots",sample,cluster,img),slide)
            write_text(20+(1-odd),"Cluster " + cluster.split("cluster_",1)[1],slide)
        
    prs.save(os.path.join(DIR,"vystupna_prezentacia.pptx"))


# =========== CLUSTER CALCULATION DBSCAN ===========
def calc_clusters(sample):
    umap = pd.read_csv(os.path.join(DIR,"Space_ranger-"+sample,"analysis","umap","gene_expression_2_components","projection.csv"))
    dbscan = DBSCAN(eps=0.75, min_samples=85)
    umap_values = umap[["UMAP-1","UMAP-2"]]
    clusters_np = dbscan.fit_predict(umap_values.values)

    clusters_df = pd.DataFrame()
    clusters_df["Barcode"] = umap["Barcode"]
    clusters_df["Cluster"] = clusters_np
    print(clusters_df['Cluster'].value_counts())
    return clusters_df

# =========== CLUSTER CALCULATION HDBSCAN ===========
def calc_hclusters(sample):
    min_cluster_sizes = range(5, 21, 2)
    num_clusters = []
    scores = []


    umap = pd.read_csv(os.path.join(DIR,"Space_ranger-"+sample,"analysis","umap","gene_expression_2_components","projection.csv"))
    umap_values = umap[["UMAP-1","UMAP-2"]]

    umap_values = umap[["UMAP-1", "UMAP-2"]].values  # čisto numerické hodnoty

    for i in min_cluster_sizes:
        clusterer = hdbscan.HDBSCAN(min_cluster_size=i)
        clusters_np = clusterer.fit_predict(umap_values)

        clusters_df = pd.DataFrame({
            "Barcode": umap["Barcode"],
            "Cluster": clusters_np
        })

        labels = clusters_df["Cluster"]
        num_clusters.append(len(labels.unique()) - 1)  # -1 pre šum

        mask = labels != -1
        if mask.sum() > 1:
            score = silhouette_score(umap_values[mask], labels[mask])
        else:
            score = 0
        scores.append(score)

    plt.figure(figsize=(10,4))
    plt.subplot(1,2,1)
    plt.plot(min_cluster_sizes, num_clusters, marker='o')
    plt.xlabel('min_cluster_size')
    plt.ylabel('Number of clusters')
    plt.show()

    plt.subplot(1,2,2)
    plt.plot(min_cluster_sizes, scores, marker='o')
    plt.xlabel('min_cluster_size')
    plt.ylabel('Silhouette score')
    plt.show()

    path = os.path.join(DIR,"Space_ranger-"+sample+".csv")
    clusters_df.to_csv(path, index=False)
    return clusters_df,path


# ============ DENSE MAP ==========
class DummyUMAP(UMAP):
    def pass_func():
        pass
    #def __init__(self, embedding,n):
     #   self.embedding_ = embedding
      #  self.n_neighbors = n

def calc_dense(sample):
    adata = sc.read_10x_h5(os.path.join(DIR,"Space_ranger-"+sample,"filtered_feature_bc_matrix.h5"))
    X = adata.X 
    
    my_clusters = calc_hclusters(sample)[0]
    clusters = my_clusters["Cluster"].to_numpy() 

    file = os.path.join(DIR,"plots","embedd-"+sample+".npy")
    if(os.path.exists(file)):
        embedding = np.load(file)
    else:
        dens_mapper = UMAP(densmap=True,dens_lambda=0.8, random_state=42).fit(X)
        np.save(file, dens_mapper.embedding_)
        embedding = dens_mapper.embedding_
    dummy = DummyUMAP()
    dummy.embedding_ = embedding
    umap.plot.points(dummy, labels=clusters, width=500, height=500)

# ============ MAIN ==============
if input_ok():
    j = 0
    create_output_dir()
    file_dic = create_dic()
    for sample in file_dic.keys():
        #calc_dense(sample)
        cluster_dic = {}
        os.makedirs(os.path.join(DIR,"plots",sample))
        #clusters_df = pd.read_csv(file_dic[sample][0])
        cluster_df,cluster_csv = calc_hclusters(sample)
        ecotypes_df = pd.read_csv(file_dic[sample][1], sep='\t')
        ecotypes_df['ID'] = ecotypes_df['ID'].str.replace('.1', '-1', regex=False)
        CLUSTERS = len(cluster_df["Cluster"].unique()) 
        for i in range(CLUSTERS):
            #cluster_dic[i+1] = filter_barcode(i+1,clusters_df)
            cluster_dic[i-1] = filter_barcode(i-1,cluster_df)
            sum_series = ecotypes_df[ecotypes_df["ID"].isin(cluster_dic[i-1])].sum()
            sum_ecotypes = sum_series[sum_series.index.isin(ECOTYPES)].sum()
            values = [sum_series[ecotype]/sum_ecotypes for ecotype in ECOTYPES]
            if not create_graphs(sample,i-1,j,cluster_csv,values): print("File problem - desired annotation not found")
        print("All charts for sample: "+sample+" have been generated")
        j+=1
    create_presentation()

